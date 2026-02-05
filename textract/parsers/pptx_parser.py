import pptx

from .utils import BaseParser


class Parser(BaseParser):
    """Extract text from pptx file using python-pptx
    """

    def extract(self, filename, **kwargs):
        presentation = pptx.Presentation(filename)
        text_runs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:  # type: ignore[attr-defined]
                    text_runs.extend(run.text for run in paragraph.runs)
        return "\n\n".join(text_runs)
